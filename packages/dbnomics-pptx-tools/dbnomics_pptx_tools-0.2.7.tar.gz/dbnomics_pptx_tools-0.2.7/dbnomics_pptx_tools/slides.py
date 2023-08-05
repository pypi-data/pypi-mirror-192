from typing import Iterable, Iterator, cast

import daiquiri
from parsy import ParseError, any_char, regex, string, whitespace
from pptx.chart.series import LineSeries, SeriesCollection
from pptx.presentation import Presentation
from pptx.shapes.graphfrm import GraphicFrame
from pptx.slide import Slide, Slides

from dbnomics_pptx_tools.charts import update_chart
from dbnomics_pptx_tools.repo import SeriesRepo
from dbnomics_pptx_tools.tables import extract_table_zones, format_table, update_table

from .metadata import PresentationMetadata, SlideMetadata

__all__ = ["delete_other_slides", "update_slides"]

logger = daiquiri.getLogger(__name__)


def delete_other_slides(prs: Presentation, *, only_slides: set[int]):
    for slide_pos, slide in enumerate(cast(Iterable[Slide], prs.slides), start=1):
        if slide_pos not in only_slides:
            delete_slide(prs, slide)


# Cf https://github.com/scanny/python-pptx/issues/67#issuecomment-296135015
def delete_slide(prs: Presentation, slide: Slide):
    slides = cast(Slides, prs.slides)
    id_dict = {slide.id: [i, slide.rId] for i, slide in enumerate(slides._sldIdLst)}
    slide_id = slide.slide_id
    prs.part.drop_rel(id_dict[slide_id][1])
    del slides._sldIdLst[id_dict[slide_id][0]]


def extract_slide_id_from_slide_notes(slide: Slide) -> str | None:
    slide_notes = slide.notes_slide.notes_text_frame.text
    parser = (
        whitespace.optional()
        >> string("slide_id", transform=str.lower)
        >> whitespace.optional()
        >> string(":")
        >> whitespace.optional()
        >> regex(r"[\w_-]+")
        << any_char.many().optional()
    )
    try:
        return parser.parse(slide_notes)
    except ParseError:
        return None


def iter_chart_shapes(slide: Slide) -> Iterator[GraphicFrame]:
    for shape in cast(list[GraphicFrame], slide.shapes):
        if shape.has_chart:
            yield shape


def iter_table_shapes(slide: Slide) -> Iterator[GraphicFrame]:
    for shape in cast(list[GraphicFrame], slide.shapes):
        if shape.has_table:
            yield shape


def iter_line_series(series_collection: SeriesCollection, *, strict: bool = True) -> Iterator[LineSeries]:
    for series in series_collection:
        if not isinstance(series, LineSeries):
            if strict:
                raise ValueError(f"Only LineSeries are expected in {series_collection!r}")
            else:
                continue
        yield series


def update_slide_charts(slide: Slide, *, repo: SeriesRepo, slide_metadata: SlideMetadata):
    slide_shapes = list(iter_chart_shapes(slide))
    for chart_pos, chart_shape in enumerate(slide_shapes, start=1):
        chart_name = chart_shape.name
        chart_full_name = f"{chart_name} ({chart_pos}/{len(slide_shapes)})"

        chart_spec = slide_metadata.charts.get(chart_name)
        if chart_spec is None:
            logger.debug(
                "Chart %s was found in the presentation file, "
                "but no corresponding block was found in the YAML file, ignoring chart",
                chart_full_name,
            )
            continue

        logger.debug("Updating chart %s...", chart_full_name)

        try:
            update_chart(
                chart_shape,
                chart_spec=chart_spec,
                repo=repo,
                slide=slide,
            )
        except Exception:
            logger.exception("Error updating chart %s, skipping", chart_full_name)


def update_slide_tables(slide: Slide, *, repo: SeriesRepo, slide_metadata: SlideMetadata):
    table_shapes = list(iter_table_shapes(slide))
    for table_pos, table_shape in enumerate(table_shapes, start=1):
        table = table_shape.table
        table_name = table_shape.name
        table_full_name = f"{table_name} ({table_pos}/{len(table_shapes)})"

        table_spec = slide_metadata.tables.get(table_name)
        if not table_spec:
            logger.debug(
                "Table %s was found in the presentation file, "
                "but no corresponding block was found in the YAML file, ignoring table, showing table preview:\n%s",
                table_full_name,
                format_table(table),
            )
            continue

        logger.debug("Updating table %s...", table_full_name)

        extract_periods = table_spec.columns is None
        table_zones = extract_table_zones(table, extract_periods=extract_periods, table_name=table_name)
        if table_zones is None:
            logger.warning(
                "Could not extract the zones of table %s, ignoring table, showing table preview:\n%s",
                table_full_name,
                format_table(table),
            )
            continue

        try:
            update_table(table, repo=repo, table_spec=table_spec, table_zones=table_zones)
        except Exception:
            logger.exception("Error updating table %s, skipping", table_full_name)


def update_slides(
    prs: Presentation,
    *,
    only_slides: set[int] | None = None,
    presentation_metadata: PresentationMetadata,
    repo: SeriesRepo,
):
    for slide_pos, slide in enumerate(cast(Iterable[Slide], prs.slides), start=1):
        if only_slides is not None and slide_pos not in only_slides:
            continue

        slide_title = slide.shapes.title.text

        slide_id = extract_slide_id_from_slide_notes(slide)
        if slide_id is None:
            logger.debug("Could not extract the slide ID from the slide notes, using the slide title %r", slide_title)

        slide_metadata = presentation_metadata.slides.get(slide_id if slide_id is not None else slide_title)
        if slide_metadata is None and slide_id is not None:
            logger.debug("Slide %d with ID %r has no metadata, trying with slide title", slide_pos, slide_id)
            slide_metadata = presentation_metadata.slides.get(slide_title)
        if slide_metadata is None:
            logger.debug(
                "Slide %d with ID %r and title %r has no metadata, ignoring slide", slide_pos, slide_id, slide_title
            )
            continue

        logger.debug("Updating slide %d with ID %r and title %r...", slide_pos, slide_id, slide_title)

        update_slide_charts(slide, repo=repo, slide_metadata=slide_metadata)

        update_slide_tables(slide, repo=repo, slide_metadata=slide_metadata)

"""Document Class for Editable and Uneditable Documents"""
# -*- coding: utf-8 -*-

# document processing
from pdf2image import convert_from_path
import docx
import pptx
from pptx.enum.shapes import MSO_SHAPE_TYPE
import imutils

# recognition stuff
import cv2
import numpy as np
from langdetect import detect
from tesserocr import PyTessBaseAPI, RIL

# basic tools
import re
import os
import shutil
import subprocess

# errors
from gtn_tools.exceptions import DocumentException
from gtn_tools.quote.helpers import ISO_to_tess
from PIL import Image

Image.MAX_IMAGE_PIXELS = 20000 * 20000
os.environ['OMP_THREAD_LIMIT'] = '1'

class Doc:
    """
    Base class for document representation.

    TODO's:
        - spelling correction

    Attributes:
    ----------
    path -- path to source file
    type -- file type
    pages -- number of pages
    text -- the (recognized) text
    source_lang -- source language of the text
    target_langss -- target languages of the text (for translation)
    words -- number of words

    Methods:
    -------
    detect_lang -- detect the language of the text
    num_words -- get the number of words
    correct_spelling -- correct the spelling of the recognized text
    postprocess -- postprocessing for recognized text
    calculate_price -- calculate a quote for the document given its attributes
    """
    def __init__(self, path: str) -> None:
        self.path = str(path)
        self.type = os.path.splitext(self.path)[1].replace('.', '').upper()

        self.pages = None
        self.text = None
        self.source_lang = None
        self.target_langs = None
        self.words = None

    def __rep__(self):
        if self.text is not None:
            return self.text
        else:
            return 'Document not yet recognized'

    # base methods
    def detect_lang(self) -> str:
        """
        Detect the language of the current document.

        Arguments:
        ---------
        source_lang -- ISO-639-1 language code of the source language
        """
        if not isinstance(self.text, str):
            raise DocumentException('The Document is not yet recognized.')
        lang = detect(self.text)
        return lang

    def num_words(self, num_incl: bool = True,
                  all_characters: bool = False) -> int:
        """
        Calculate the number of words in the document.

        Arguments:
        ---------
        num_incl -- consider the numbers and digits as words
        all_characters -- consider all characters as words
        """
        if not isinstance(self.text, str):
            raise DocumentException('The Document is not yet recognized.')
        word_list = self.text.split(' ')
        if not all_characters:
            if num_incl:
                regex = re.compile('[^a-zA-Z0-9]')
            else:
                regex = re.compile('[^a-zA-Z]')
            for i in range(len(word_list)):
                word_list[i] = regex.sub('', word_list[i])
        while '' in word_list:
            word_list.remove('')
        return len(word_list)

    def correct_spelling(self) -> None:
        """
        Correct the spelling of the recognized Text.
        """
        # TODO
        return None

    def postprocess(self, target_langs: str, source_lang: str = None) -> None:
        """
        Postprocess recognized text. Postprocessing includes detecting the
        source language if not given as well as calculating the number of words
        in the document.
        TODO add spelling correction to postprocessing.

        Arguments:
        ---------
        source_lang -- ISO-639-1 language code of the source language. Supply
        `None` to automatically detect the source_lang
        """
        if not isinstance(self.text, str):
            raise DocumentException('The Document is not yet recognized.')
        if source_lang is None:
            self.source_lang = self.detect_lang(source_lang)
        else:
            self.source_lang = source_lang
        self.target_langs = target_langs
        self.words = self.num_words()
        return None

    def document_workflow(self, source_lang: str, target_langs: str,
                          prc_dir: str, uneditable_files: list = None):
        """
        Starts the recognition workflow. If the filetype is in the
        'uneditale_files'-list, the document will be converted to jpg
        and preprocessed for ocr. Then the jpg files will be recognized and
        boxes will be drawn around every recognized word. The recognized
        text and the jpg-files with boxes will be saved to 'prc_dir'. A light
        postprocessing will be applied.

        Arguments:
        ---------
        target_langs: The target language of the doc
        source_lang: The source language of the doc
        prc_dir: The directory to save the output to
        uneditable_files: A list of allowed uneditable files. Will be PDF,
        JPG if nothing provided
        """
        if not os.path.isdir(prc_dir):
            os.makedirs(prc_dir)
        if uneditable_files is None:
            uneditable_files = ['PDF', 'JPG']
        # Step 1: convert pdf to jpeg
        uneditable = False
        if self.type in uneditable_files:
            uneditable = True
        if uneditable:
            self.convert(prc_dir)
        if uneditable:
            # Step 2: preprocess images/editable documents
            self.preprocess()
            # Step 3: print boxes around words
            self.print_boxes(config=f'-l {ISO_to_tess(source_lang)}')
            # Step 4: recognize text in document
            docname = (self.path[self.path.rfind('/')
                       + 1:self.path.rfind('.')])
            self.recognize(output_file=os.path.join(prc_dir, docname +
                                                    '_recognized_text.txt'),
                           config=f'-l {ISO_to_tess(source_lang)}')
        else:
            # Skip steps 2 and 3 if file is editable
            self.recognize()
        # Step 5: postprocess document
        self.postprocess(target_langs, source_lang)

        return None


class Uneditable(Doc):
    """
    Class for uneditable file formats. Currently, pdf and jpg filetypes are
    supported.

    General Doc Attributes:
    ----------------------
    path -- path to source file
    type -- file type
    pages -- number of pages
    text -- the (recognized) text
    source_lang -- source language of the text
    target_langs -- target language of the text (for translation)
    words -- number of words

    Uneditable Attributes:
    ---------------------
    page_paths -- list of paths for image files
    img -- list of images represented as PIL.Image
    cv2_img -- list of images represented as cv2.Image
    preprocessed -- list of preprocessed files
    rotated -- list of rotated image files
    preprocessing_methods -- the methods applied to the document during
    preprocessing
    config -- custom config for tesseract

    General Doc Methods:
    -------------------
    detect_lang -- detect the language of the text
    num_words -- get the number of words
    correct_spelling -- correct the spelling of the recognized text
    postprocess -- postprocessing for recognized text
    calculate_price -- calculate a quote for the document given its attributes

    Methods:
    -------
    ("do-methods") -- helper methods for preprocessing
    set_page -- set pages manually as JPEG
    convert -- convert single PDF files to a JPEG file per page
    print_boxes -- print boxes around the words
    preprocess -- preprocess files
    recognize -- recognize files
    """
    def __init__(self, path: str) -> None:
        super().__init__(path)
        self.page_paths = []
        self.preprocessing_methods = None
        self.config = None
        self.api = PyTessBaseAPI()
        if self.type == 'JPEG' or self.type == 'JPG' or self.type == 'PNG':
            self.page_paths.append(path)
        return None

    # preprocessing methods
    def do_grayscale(self, image: np.ndarray) -> np.ndarray:
        """
        Grayscale Image

        Arguments:
        ---------
        images -- list of cv2.Images
        """
        return cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)

    def do_remove_noise(self, image: np.ndarray) -> np.ndarray:
        """
        Remove noise

        Arguments:
        ---------
        images -- list of cv2.Images
        """
        return cv2.medianBlur(image, 5)

    def do_threshold(self, image: np.ndarray) -> np.ndarray:
        """
        Threshold Image

        Arguments:
        ---------
        images -- list of cv2.Images
        """
        return cv2.threshold(image, 0, 255,
                             cv2.THRESH_BINARY + cv2.THRESH_OTSU)[1]

    def do_dilate(self, image: np.ndarray) -> np.ndarray:
        """
        Dilate Image

        Arguments:
        ---------
        images -- list of cv2.Images
        """
        kernel = np.ones((5, 5), np.uint8)
        return cv2.dilate(image, kernel, iterations=1)

    def do_erode(self, image: np.ndarray) -> np.ndarray:
        """
        Erode Image

        Arguments:
        ---------
        images -- list of cv2.Images
        """
        kernel = np.ones((5, 5), np.uint8)
        return cv2.erode(image, kernel, iterations=1)

    def do_opening(self, image: np.ndarray) -> np.ndarray:
        """
        Opening Image

        Arguments:
        ---------
        images -- list of cv2.Images
        """
        kernel = np.ones((5, 5), np.uint8)
        return cv2.morphologyEx(image, cv2.MORPH_OPEN, kernel)

    def do_rotate(self, image: np.ndarray) -> np.ndarray:
        """
        Rotate Image

        Arguments:
        ---------
        images -- list of cv2.Images
        """
        orientation = self.api.DetectOS()
        if orientation is not None:
            angle = self.api.DetectOS().get('orientation')
        else:
            angle = 0
        return imutils.rotate_bound(image, angle)

    # methods for text recognition
    def set_page(self, path: str) -> None:
        """
        Set the pages directly if input file is JPG

        Arguments:
        ---------
        path -- path to jpg file
        """
        self.page_paths.append(path)
        return None

    def convert(
        self,
        out_dir: str,
        dpi: int = 300,
        fmt: str = 'jpg',
        threads: int = 1
    ) -> None:
        """
        Convert PDF to editable format.

        Arguments:
        ---------
        save -- Should the image file(s) be saved? Default is False.
        path -- The directory to save the image file(s).
        Default is the same as the directory of the original file
        dpi -- Quality of the image(s)
        transparent -- Background transparent instead of white
        fmt -- The format to save the image file(s)
        """
        if self.type == 'PDF':
            self.page_paths = convert_from_path(self.path, dpi=dpi, fmt=fmt,
                                                output_folder=out_dir,
                                                thread_count=threads,
                                                paths_only=True)
        elif self.type in ["JPG", "JPEG"]:
            for i in range(len(self.page_paths)):
                d = self.page_paths[i]
                s = self.page_paths[i].replace("/src/", "/prc/")
                shutil.copyfile(d, s)
                self.page_paths[i] = s
        elif self.type in ["PNG"]:
            for i in range(len(self.page_paths)):
                with Image.open(self.page_paths[i]) as im:
                    if im.mode != 'RGB':
                        im = im.convert('RGB')
                    filename = self.page_paths[i].replace(".png", ".jpg").replace("/src/", "/prc/")
                    im.save(filename, format("JPEG"))
                    self.page_paths[i] = filename
        self.pages = len(self.page_paths)
        return None

    def preprocess(self, method: str = None) -> None:
        """
        Preprocess documents for recognition. If no method is supplied, the
        default workflow is 1. rotate, 2. grayscale, 3. threshold. Else, the
        supplied method will be applied. For applying multiple methods, call
        the function multiple times.

        Arguments:
        ---------
        method -- method to be applied for preprocessing. Choices: grayscale,
        remove_noise, threshold, dilate, erode, opening, rotate
        """
        if method is not None:
            if method not in ['grayscale', 'remove_noise', 'threshold',
                              'dilate', 'erode', 'opening', 'rotate']:
                raise DocumentException('Invalid method supplied.')
            if self.preprocessing_methods != '':
                self.preprocessing_methods += method + ', '
            else:
                self.preprocessing_methods = method + ', '
            images = self.cv2_img if self.preprocessed is None \
                else self.preprocessed
            fun = str('do_' + method)
            self.preprocessed = getattr(self, fun)(images)
            return None
        else:
            self.preprocessing_methods = 'rotate, grayscale, thresholding'
            # Step 1: rotate
            for impath in self.page_paths:
                image = cv2.imread(impath)
                self.api.SetImageFile(impath)
                # Step 1: rotate
                rotated = self.do_rotate(image)
                # Step 2: grayscale
                gray = cv2.cvtColor(rotated, cv2.COLOR_BGR2GRAY)
                # Step 3: thresholding
                thresh = cv2.threshold(gray, 0, 255,
                                       cv2.THRESH_BINARY | cv2.THRESH_OTSU)[1]

                cv2.imwrite(impath, thresh)
            return None

    def print_boxes(self, config: str = None) -> None:
        """
        Prints boxes around words in preprocessed documents. If no config is
        provided the method looks for the global config. If no global config
        is defined the method fails.

        Arguments:
        ---------
        out_dir -- the directory to save the jpg files
        config -- tesseract config, incl. the language configuration
        """
        if config is None:
            config = self.config
        try:
            self.api = PyTessBaseAPI(lang=config.split('-l ', 1)[1])
        except:
            self.api = PyTessBaseAPI(lang='eng')

        for impath in self.page_paths:
            image = cv2.imread(impath)
            self.api.SetImageFile(impath)
            boxes = self.api.GetComponentImages(RIL.WORD, True)
            for i in range(len(boxes)):
                box = boxes[i][1]
                x, y, w, h = box['x'], box['y'], box['w'], box['h']
                cv2.rectangle(image, (x, y), (x + w, y + h), color=(255, 0, 255), thickness=3)
            cv2.imwrite(impath, image)
        return None

    def recognize(self, output_file: str = None, config: str = None):
        """
        Recognize the document's text via tesseract-ocr.

        Arguments:
        ---------
        output_file -- optional path to textfile where the output should be
        saved
        config -- configuration for the tesseract command line tool must
        contain at least the language of the document to be recognized
        (e.g. `-l deu`)
        """
        if config is None:
            config = self.config
        else:
            self.config = config
        if output_file is not None:
            f = open(output_file, 'a')
        try:
            self.api = PyTessBaseAPI(lang=config.split('-l ', 1)[1])
        except:
            self.api = PyTessBaseAPI(lang='eng')
        output = ''
        for impath in self.page_paths:
            self.api.SetImageFile(impath)
            text = self.api.GetUTF8Text()
            text = text.replace('-\n', '')
            text = text.replace('\n', ' ')
            output += str(' ' + text)
            if output_file is not None:
                f.write(text)
        if output_file is not None:
            f.close
        self.text = output
        if output_file is None:
            return self.text
        else:
            return None


class Editable(Doc):
    """
    Class for editable document types. Currently, the following extensions are
    supported: doc, docx, ppt, pptx.

    General Doc Attributes:
    ----------------------
    path -- path to source file
    type -- file type
    pages -- number of pages
    text -- the (recognized) text
    source_lang -- source language of the text
    target_langs -- target language of the text (for translation)
    words -- number of words

    Attributes:
    ----------
    None

    General Doc Methods:
    -------------------
    detect_lang -- detect the language of the text
    num_words -- get the number of words
    correct_spelling -- correct the spelling of the recognized text
    postprocess -- postprocessing for recognized text
    calculate_price -- calculate a quote for the document given its attributes

    Methods:
    -------
    ("do-methods") -- helper methods for preprocessing
    recognize -- recognize text in files
    """
    def __init__(self, path: str) -> None:
        super().__init__(path)

    def do_get_shapes(self, prs):
        """
        Helper for pptx processing. Returns a list of non-group shapes and a
        list of group shapes for further processing

        Arguments:
        ---------
        prs -- a pptx.Presentation to extract shapes
        """
        shapes = []
        groups = []
        for slide in prs.slides:
            for shape in slide.shapes:
                shapes.append(shape)
        shape_types = [shape.shape_type for shape in shapes]
        for type in shape_types:
            if type == MSO_SHAPE_TYPE.GROUP:
                groups.append(shapes.pop(shape_types.index(type)))
                shape_types = [shape.shape_type for shape in shapes]
        return shapes, groups

    def do_extract_children(self, groups: list):
        """
        Helper for pptx processing. Extract all (nested) children from group
        shapes.

        Arguments:
        ---------
        groups -- a list of group shapes as returned by self.do_get_shapes()
        """
        final_children = []
        group_children = []
        for group in groups:
            for shape in group.shapes:
                if not shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    final_children.append(shape)
                else:
                    group_children.append(shape)
        return final_children, group_children

    def do_convert_ppt(self) -> None:
        """
        Convert ppt to pptx format to be recognized.
        """
        subprocess.call(['unoconv', '-d', 'document', '--format=pptx',
                         self.path])
        self.path = self.path + 'x'
        self.type = self.path[self.path.rfind('.')+1:].upper()
        return None

    def do_process_pptx(self):
        """
        Recognize text in pptx files and return the number of words or
        unprocessed text.

        Arguments:
        ---------
        text -- if True, the recognized text, if False the number of words will
        be returned
        """
        prs = pptx.Presentation(self.path)
        shapes, groups = self.do_get_shapes(prs)
        while True:
            final_children, group_children = self.do_extract_children(groups)
            shapes.extend(final_children)
            groups = group_children
            if group_children == []:
                break
        textlist = []
        for shape in shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                textlist.append(paragraph.text)
        return str(' '.join(textlist))

    def do_convert_doc(self) -> None:
        """
        Convert doc to docx format to be recognized.
        """
        subprocess.call(['unoconv', '-d', 'document', '--format=docx',
                         self.path])
        self.path = self.path + 'x'
        self.type = self.path[self.path.rfind('.')+1:].upper()
        return None

    def do_process_docx(self):
        """
        Recognize text in docx files and return the number of words or
        unprocessed text.

        Arguments:
        ---------
        text -- if True, the recognized text, if False the number of words will
        be returned
        """
        document = docx.Document(self.path)
        textlist = []
        for paratext in document.paragraphs:
            textlist.append(paratext.text)
        return str(' '.join(textlist))

    def recognize(self) -> None:
        """
        Recognize the text in editable files.
        """
        if self.type.upper() == 'DOC':
            self.do_convert_doc()
        elif self.type.upper() == 'PPT':
            self.do_convert_ppt()
        if self.type.upper() == "DOCX":
            self.text = self.do_process_docx()
        elif self.type.upper() == "PPTX":
            self.text = self.do_process_pptx()
        return None
